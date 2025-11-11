import os
import tempfile
from typing import List, Optional

import streamlit as st
from sentence_transformers import SentenceTransformer
import numpy as np

# Document reading
from io import BytesIO
from pdfminer.high_level import extract_text as pdf_extract_text
import docx
from pptx import Presentation
import pandas as pd

# Vector store
import faiss

# OpenAI (optional) for generative answers
import openai

# ---------- Helpers for file reading (unchanged) ----------

def read_pdf(file_bytes: bytes) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp.flush()
        text = pdf_extract_text(tmp.name)
    return text


def read_docx(file_bytes: bytes) -> str:
    doc = docx.Document(BytesIO(file_bytes))
    full = []
    for para in doc.paragraphs:
        if para.text:
            full.append(para.text)
    return "\n".join(full)


def read_txt(file_bytes: bytes) -> str:
    return file_bytes.decode(errors="ignore")


def read_pptx(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    slides_text = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slides_text.append(f"[Slide {slide_no}] {shape.text}")
    return "\n".join(slides_text)


def read_excel(file_bytes: bytes) -> str:
    excel_data = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
    text_parts = []
    for sheet_name, df in excel_data.items():
        text_parts.append(f"--- Sheet: {sheet_name} ---")
        try:
            text_parts.append(df.to_string(index=False))
        except Exception:
            rows = []
            for r in df.values:
                rows.append("\t".join([str(c) for c in r]))
            text_parts.append("\n".join(rows))
    return "\n".join(text_parts)


def split_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    tokens = text.split()
    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunk = " ".join(tokens[start:end])
        chunks.append(chunk)
        start = end - overlap
        if start < 0:
            start = 0
    return chunks


def embed_texts(model, texts: List[str]) -> np.ndarray:
    embs = model.encode(texts, show_progress_bar=False, convert_to_numpy=True)
    norms = np.linalg.norm(embs, axis=1, keepdims=True)
    norms[norms == 0] = 1
    embs = embs / norms
    return embs

# ---------- Cloud provider helpers (OpenAI + Gemini) ----------

def get_openai_embeddings(api_key: str, texts: List[str], model="text-embedding-3-small") -> np.ndarray:
    if not api_key:
        raise ValueError("OpenAI API key not provided.")
    openai.api_key = api_key
    # Note: for large number of chunks consider batching to avoid request limits
    resp = openai.Embedding.create(input=texts, model=model)
    embs = np.array([np.array(r["embedding"]) for r in resp["data"]])
    norms = np.linalg.norm(embs, axis=1, keepdims=True)
    norms[norms == 0] = 1
    embs = embs / norms
    return embs

def openai_generate(api_key: str, prompt: str, model="gpt-4o-mini") -> str:
    if not api_key:
        raise ValueError("OpenAI API key not provided.")
    openai.api_key = api_key
    completion = openai.ChatCompletion.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=512,
        temperature=0.0,
    )
    return completion["choices"][0]["message"]["content"]


def get_gemini_embeddings(api_key: str, texts: List[str], model="embed-gecko-001") -> np.ndarray:
    """
    Attempt to use the google-generative-ai library to get embeddings from Gemini/Gecko.
    Requires: pip install google-generative-ai
    """
    try:
        import google.generativeai as genai  # type: ignore
    except Exception as e:
        raise RuntimeError("google-generative-ai package not installed. Install with `pip install google-generative-ai`.") from e

    if not api_key:
        raise ValueError("Gemini API key not provided.")
    genai.configure(api_key=api_key)

    # The exact shape of the response may vary by library version; adapt if needed.
    resp = genai.embeddings.create(model=model, input=texts)
    # response usually: {'data': [{'embedding': [...]} , ...]}
    embs = np.array([np.array(item["embedding"]) for item in resp["data"]])
    norms = np.linalg.norm(embs, axis=1, keepdims=True)
    norms[norms == 0] = 1
    embs = embs / norms
    return embs

def gemini_generate(api_key: str, prompt: str, model="gemini-1.0") -> str:
    """
    Attempt to create a chat completion via google-generative-ai.
    The exact API surface could differ; adapt model name as needed.
    """
    try:
        import google.generativeai as genai  # type: ignore
    except Exception as e:
        raise RuntimeError("google-generative-ai package not installed. Install with `pip install google-generative-ai`.") from e

    if not api_key:
        raise ValueError("Gemini API key not provided.")
    genai.configure(api_key=api_key)

    # The chat API: genai.chat.create or genai.responses.create depending on version.
    # We'll try genai.chat.create first, then genai.generate if needed.
    try:
        resp = genai.chat.create(model=model, messages=[{"role": "user", "content": prompt}])
        # Many versions return text in resp.last or resp["candidates"][0]["content"]
        if hasattr(resp, "last"):
            return resp.last
        # fallback:
        return resp["candidates"][0]["content"]
    except Exception:
        # fallback to responses API
        resp2 = genai.responses.create(model=model, input=prompt)
        # response text might be in resp2.output[0].content[0].text or similar; adapt safely:
        # We attempt to extract a plain text string:
        try:
            # try common structure
            return resp2.output[0].content[0]["text"]
        except Exception:
            # last resort: stringify
            return str(resp2)

# ---------- Streamlit UI ----------

st.set_page_config(page_title="RAG Streamlit Demo (OpenAI + Gemini)", layout="wide")
st.title("RAG-based QA Agent — Streamlit (OpenAI / Gemini support)")

st.sidebar.header("Settings")

embedding_choice = st.sidebar.selectbox(
    "Embedding model",
    (
        "sentence-transformers/all-MiniLM-L6-v2",
        "sentence-transformers/all-mpnet-base-v2",
        "openai (text-embedding-3-small)",
        "gemini (embed-gecko-001)",
    ),
)

# provider for generation (if you want separate choice)
generation_provider = st.sidebar.selectbox(
    "Generation provider (LLM)",
    (
        "none",
        "openai",
        "gemini",
    ),
    index=0,
)

# API keys
openai_api_key = ""
gemini_api_key = ""
# Only show key fields when relevant
if embedding_choice.startswith("openai") or generation_provider == "openai":
    openai_api_key = st.sidebar.text_input("OpenAI API Key", type="password")
if embedding_choice.startswith("gemini") or generation_provider == "gemini":
    gemini_api_key = st.sidebar.text_input("Gemini API Key", type="password (Google)")

chunk_size = st.sidebar.number_input("Chunk size (words approx)", value=500, step=100)
overlap = st.sidebar.number_input("Chunk overlap (words)", value=50, step=10)

st.sidebar.markdown("---")
store_persist = st.sidebar.checkbox("Persist FAISS index to disk (./faiss_index)")

# In-memory store
index = None
doc_texts: List[str] = []

# Load local sentence-transformer models (as before)
@st.cache_resource
def load_sentence_model(name: str, local_dir: Optional[str] = None):
    if local_dir and os.path.isdir(local_dir):
        st.info(f"Loading SentenceTransformer from local folder: {local_dir}")
        return SentenceTransformer(local_dir)
    st.info(f"Loading SentenceTransformer from hub: {name} (may download on first run)")
    return SentenceTransformer(name)

# map sidebar option -> local folders (adjust to your paths)
local_model_map = {
    "sentence-transformers/all-MiniLM-L6-v2": r"C:\Users\m95071_adm\PyCharmMiscProject\myproject\all-MiniLM-L6-v2",
    "sentence-transformers/all-mpnet-base-v2": r"C:\Users\m95071_adm\PyCharmMiscProject\myproject\all-mpnet-base-v2",
}

embedding_model = None
if embedding_choice.startswith("sentence-transformers"):
    local_dir = local_model_map.get(embedding_choice)
    try:
        embedding_model = load_sentence_model(embedding_choice, local_dir=local_dir)
        st.success(f"Embedding model ready: {embedding_choice}")
    except Exception as e:
        st.error(f"Failed to load embedding model '{embedding_choice}': {e}")
        embedding_model = None

# Upload UI
st.header("Upload file")
uploaded = st.file_uploader("Upload PDF, DOCX, PPTX, TXT or Excel", type=["pdf", "docx", "pptx", "txt", "xlsx", "xls"])

if uploaded is not None:
    file_bytes = uploaded.read()
    st.info(f"Processing {uploaded.name} — this may take a few seconds...")
    name = uploaded.name.lower()
    if name.endswith(".pdf"):
        text = read_pdf(file_bytes)
    elif name.endswith(".docx"):
        text = read_docx(file_bytes)
    elif name.endswith(".pptx"):
        text = read_pptx(file_bytes)
    elif name.endswith(".xlsx") or name.endswith(".xls"):
        try:
            text = read_excel(file_bytes)
        except Exception as e:
            st.error(f"Failed to read Excel file: {e}")
            st.stop()
    else:
        text = read_txt(file_bytes)

    st.success(f"Read {len(text.split())} words from file.")

    # chunk
    chunks = split_text(text, chunk_size=chunk_size, overlap=overlap)
    st.write(f"Created {len(chunks)} chunks (chunk_size={chunk_size}, overlap={overlap}).")

    # embed selection branching
    if embedding_choice.startswith("sentence-transformers"):
        if embedding_model is None:
            st.error("Embedding model is not loaded. Please check model path or internet connection.")
            st.stop()
        embs = embed_texts(embedding_model, chunks)
    elif embedding_choice.startswith("openai"):
        try:
            embs = get_openai_embeddings(openai_api_key, chunks)
        except Exception as e:
            st.error(f"OpenAI embedding error: {e}")
            st.stop()
    elif embedding_choice.startswith("gemini"):
        try:
            embs = get_gemini_embeddings(gemini_api_key, chunks)
        except Exception as e:
            st.error(f"Gemini embedding error: {e}")
            st.stop()
    else:
        st.error("Unknown embedding choice")
        st.stop()

    embs = embs.astype("float32")
    dim = embs.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embs)
    doc_texts = chunks

    if store_persist:
        os.makedirs("./faiss_index", exist_ok=True)
        faiss.write_index(index, "./faiss_index/index.faiss")
        with open("./faiss_index/texts.txt", "w", encoding="utf-8") as f:
            for c in doc_texts:
                f.write(c.replace("\n", " ") + "\n<<CHUNK>>\n")
        st.success("Persisted FAISS index and texts to ./faiss_index/")

    st.session_state["index_ready"] = True
    st.session_state["doc_texts"] = doc_texts
    st.session_state["faiss_index"] = index

# Querying
st.header("Ask a question about the uploaded file")
query = st.text_input("Question")
if st.button("Get Answer"):
    if "index_ready" not in st.session_state:
        st.error("Upload and process a file first.")
    elif not query:
        st.error("Please enter a question.")
    else:
        idx = st.session_state["faiss_index"]
        texts = st.session_state["doc_texts"]

        # embed query
        if embedding_choice.startswith("sentence-transformers"):
            if embedding_model is None:
                st.error("Embedding model is not loaded. Cannot embed query.")
                st.stop()
            q_emb = embedding_model.encode([query], convert_to_numpy=True)
            q_emb = q_emb / np.linalg.norm(q_emb, axis=1, keepdims=True)
        elif embedding_choice.startswith("openai"):
            try:
                q_emb = get_openai_embeddings(openai_api_key, [query])
            except Exception as e:
                st.error(f"OpenAI query embedding error: {e}")
                st.stop()
        elif embedding_choice.startswith("gemini"):
            try:
                q_emb = get_gemini_embeddings(gemini_api_key, [query])
            except Exception as e:
                st.error(f"Gemini query embedding error: {e}")
                st.stop()
        else:
            st.error("Unknown embedding choice for query.")
            st.stop()

        # search
        k = 5
        D, I = idx.search(q_emb.astype("float32"), k)
        I = I[0]
        D = D[0]
        retrieved = [texts[i] for i in I if i != -1]

        st.subheader("Top retrieved chunks (context)")
        for score, chunk in zip(D, retrieved):
            st.write(f"Score: {float(score):.4f}")
            st.write(chunk[:1000] + ("..." if len(chunk) > 1000 else ""))
            st.markdown("---")

        # Generation step (choose provider)
        if generation_provider == "openai":
            if not openai_api_key:
                st.error("OpenAI API key required for generation.")
                st.stop()
            prompt = (
                "You are an assistant that answers questions using the provided context.\n"
                "If the answer isn't in the context, say 'I don't know based on the provided document.'\n\n"
                "Context:\n" + "\n---\n".join(retrieved) + f"\n\nQuestion: {query}\nAnswer:"
            )
            try:
                ans = openai_generate(openai_api_key, prompt)
                st.subheader("Answer (generated - OpenAI)")
                st.write(ans)
            except Exception as e:
                st.error(f"OpenAI generation error: {e}")
        elif generation_provider == "gemini":
            if not gemini_api_key:
                st.error("Gemini API key required for generation.")
                st.stop()
            prompt = (
                "You are an assistant that answers questions using the provided context.\n"
                "If the answer isn't in the context, say 'I don't know based on the provided document.'\n\n"
                "Context:\n" + "\n---\n".join(retrieved) + f"\n\nQuestion: {query}\nAnswer:"
            )
            try:
                ans = gemini_generate(gemini_api_key, prompt)
                st.subheader("Answer (generated - Gemini)")
                st.write(ans)
            except Exception as e:
                st.error(f"Gemini generation error: {e}")
        else:
            st.subheader("Extractive answer — top contexts shown above")
            st.info("Enable a generation provider (OpenAI or Gemini) in the sidebar to produce a synthesized answer.")

st.markdown("---")
st.write(
    "Notes: This demo supports local sentence-transformers models, OpenAI, and Gemini for embeddings and generation. "
    "For Gemini support install `google-generative-ai` and provide a valid API key. "
    "Model names used for Gemini (embed-gecko-001 / gemini-1.0) may need to be adjusted to your account and SDK version."
)
